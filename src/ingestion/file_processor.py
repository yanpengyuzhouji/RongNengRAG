"""
文件处理器 — 模块化入库核心
接收单个或多个文件路径，执行完整的 parse → chunk → embed → insert 管道

设计原则:
  - 每个文件独立处理，可随时添加/删除/重建索引
  - 不依赖目录扫描，完全由调用方驱动
  - 返回详细的处理报告，便于上层 (API/UI/CLI) 展示进度
"""

import os
import sys
import time
import json
import hashlib
import sqlite3
import shutil
import threading
from pathlib import Path
from datetime import datetime
from typing import Optional, List, Dict, Callable
from dataclasses import dataclass, field
from enum import Enum

import yaml

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from ingestion.pdf_parser import PDFParser
from ingestion.chunker import Chunker, Chunk
from ingestion.embedder import Embedder, create_text_for_embedding
from ingestion.milvus_store import MilvusStore
from ingestion.metadata_sync import (
    inherit_registry_metadata,
    normalize_metadata_updates,
    synchronize_metadata,
)
from ingestion.ocr_fallback import choose_page_text
from ingestion.document_editor import (
    LayoutEditError,
    LayoutRevisionConflict,
    apply_layout_edits,
    count_visual_assets,
    layout_page_texts,
    layout_revision,
)


class FileStatus(str, Enum):
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"
    DELETED = "deleted"


@dataclass
class ProcessResult:
    """单个文件的处理结果"""
    file_path: str
    file_hash: str
    file_name: str
    status: FileStatus
    chunks_created: int = 0
    chars_extracted: int = 0
    parse_time_ms: float = 0
    embed_time_ms: float = 0
    total_time_ms: float = 0
    error_message: str = ""
    # 元数据
    domain: str = ""
    category: str = ""
    doc_number: str = ""
    file_type: str = ""


@dataclass
class BatchResult:
    """批量处理结果"""
    total: int
    success: int
    failed: int
    results: List[ProcessResult] = field(default_factory=list)
    total_time_ms: float = 0

    @property
    def success_rate(self) -> float:
        return self.success / self.total if self.total > 0 else 0


class FileProcessor:
    """
    模块化文件处理器

    用法:
        processor = FileProcessor()
        result = processor.process("D:/path/to/file.pdf")
        batch_result = processor.process_batch(["file1.pdf", "file2.pdf"])
        processor.delete("file_hash_or_path")
        processor.reindex("file_hash_or_path")
    """

    def __init__(self, config_path: str = None):
        from config import load_config, get_config_path, ensure_data_dirs
        self.config = load_config(config_path)
        self.config_path = config_path or get_config_path()

        ocr_cfg = self.config.get("ocr", {})
        self.pdf_parser = PDFParser(
            min_text_chars=ocr_cfg.get("min_text_chars", 50),
            ocr_config=ocr_cfg,
        )
        self.chunker = Chunker(config_path)
        self.embedder = None  # 延迟加载
        self.store = MilvusStore(config_path)
        self._metadata_lock = threading.RLock()
        self._index_lock = threading.RLock()

        # 数据库路径 (已由 load_config 解析为绝对路径)
        self.db_path = self.config["paths"]["metadata_db"]
        self.uploads_dir = Path(self.config["paths"]["uploads_dir"])
        self.uploads_dir.mkdir(parents=True, exist_ok=True)

        self._init_registry()

    def _chunk_text_safe(self, text: str, file_meta: dict) -> list:
        """分块文本 — 超大文本分批处理防止 chunker 递归溢出"""
        if len(text) <= 100000:
            return self.chunker.chunk_text_document(text, file_meta)

        # 分批: 按双换行拆段，单批 ≤30000 字符
        sections = [s.strip() for s in text.split("\n\n") if s.strip()]
        if len(sections) == 1 and len(sections[0]) > 30000:
            sections = [s.strip() for s in text.split("\n") if s.strip()]

        all_chunks = []
        batch = ""
        for sec in sections:
            if len(batch) + len(sec) < 30000:
                batch = (batch + "\n\n" + sec) if batch else sec
            else:
                if batch:
                    all_chunks.extend(self._chunk_batch_or_fallback(batch, file_meta))
                batch = sec
        if batch:
            all_chunks.extend(self._chunk_batch_or_fallback(batch, file_meta))
        return all_chunks

    def _chunk_batch_or_fallback(self, batch: str, file_meta: dict) -> list:
        """尝试分块一批文本，MemoryError 时强制按 2000 字符拆分"""
        try:
            return self.chunker.chunk_text_document(batch, file_meta)
        except MemoryError:
            chunks = []
            for i in range(0, len(batch), 2000):
                sub = batch[i:i + 2000].strip()
                if sub:
                    try:
                        chunks.extend(self.chunker.chunk_text_document(sub, file_meta))
                    except Exception:
                        pass
            return chunks

    def _get_db_connection(self):
        """获取 SQLite 连接，统一启用 WAL 模式 + busy_timeout，避免 database is locked"""
        conn = sqlite3.connect(self.db_path, timeout=30)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA busy_timeout=5000")
        return conn

    def _init_registry(self):
        """初始化文件注册表 (SQLite)"""
        conn = self._get_db_connection()
        cursor = conn.cursor()

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS file_registry (
                file_hash TEXT PRIMARY KEY,
                original_path TEXT NOT NULL,
                stored_path TEXT,
                file_name TEXT NOT NULL,
                file_size INTEGER,
                file_type TEXT,
                status TEXT DEFAULT 'pending',
                chunks_count INTEGER DEFAULT 0,
                chars_count INTEGER DEFAULT 0,
                domain TEXT,
                category TEXT,
                doc_number TEXT,
                error_message TEXT,
                parse_time_ms REAL DEFAULT 0,
                embed_time_ms REAL DEFAULT 0,
                created_at TEXT DEFAULT (datetime('now', 'localtime')),
                updated_at TEXT DEFAULT (datetime('now', 'localtime')),
                reindex_count INTEGER DEFAULT 0
            )
        """)

        cursor.execute("""
            CREATE INDEX IF NOT EXISTS idx_registry_status ON file_registry(status)
        """)
        cursor.execute("""
            CREATE INDEX IF NOT EXISTS idx_registry_domain ON file_registry(domain)
        """)

        conn.commit()
        conn.close()

        self._recover_stuck_files()

    def _recover_stuck_files(self):
        """启动时恢复卡在 processing 状态的文件 (上次崩溃残留)"""
        conn = self._get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            "UPDATE file_registry SET status='pending', error_message='' "
            "WHERE status='processing'"
        )
        n = cursor.rowcount
        if n > 0:
            print(f"[recover] 恢复 {n} 个卡在 processing 状态的文件 -> pending")
        conn.commit()
        conn.close()

    def compute_hash(self, filepath: str) -> str:
        """计算文件 SHA256"""
        sha256 = hashlib.sha256()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def process(self, file_path: str,
                domain: str = None,
                category: str = None,
                progress_callback: Callable[[str, float], None] = None,
                force_reindex: bool = False,
                ) -> ProcessResult:
        """
        处理单个文件: parse → chunk → embed → insert

        Args:
            file_path: 文件路径 (绝对路径)
            domain: 手动指定专业域 (可选，默认自动推断)
            category: 手动指定文档类目 (可选)
            progress_callback: 进度回调 (阶段名, 0~1进度)

        Returns:
            ProcessResult
        """
        from ingestion.file_walker import FileWalker

        if not os.path.exists(file_path):
            return ProcessResult(
                file_path=file_path, file_hash="", file_name=Path(file_path).name,
                status=FileStatus.FAILED, error_message=f"文件不存在: {file_path}"
            )

        t_start = time.time()
        file_path_obj = Path(file_path)
        file_name = file_path_obj.name
        file_hash = self.compute_hash(file_path)
        file_ext = file_path_obj.suffix.lower()
        file_size = file_path_obj.stat().st_size

        # 检查是否已入库
        existing = self._get_registry(file_hash)
        if existing and existing["status"] == "completed" and not force_reindex:
            # 旧版本只保存了文本 chunks，没有 layout cache。此类文件即使
            # 已完成，也必须重新走 8001 pipeline 才能补齐版面数据。
            layout_cache = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache")) / f"{file_hash}.layout.json"
            needs_layout_rebuild = (
                file_ext in (".pdf", ".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".gif")
                and self.config.get("ocr", {}).get("enabled", False)
                and not layout_cache.exists()
            )
            if needs_layout_rebuild:
                print(f"[OCR] 已完成文件缺少版面缓存，重新调用 8001 pipeline: {file_name}", flush=True)
                force_reindex = True
            else:
                # 如果文件名变了（重命名再上传），更新注册表中的 file_name 和 original_path
                old_name = existing.get("file_name", "")
                if file_name != old_name:
                    self._upsert_registry(file_hash, file_name=file_name,
                                          original_path=file_path)
                return ProcessResult(
                    file_path=file_path, file_hash=file_hash, file_name=file_name,
                    status=FileStatus.COMPLETED,
                    chunks_created=existing["chunks_count"],
                    chars_extracted=existing["chars_count"],
                    domain=existing.get("domain", ""),
                    error_message="文件已入库，无需重复处理"
                )

        result = ProcessResult(
            file_path=file_path, file_hash=file_hash, file_name=file_name,
            status=FileStatus.PROCESSING, file_type=file_ext
        )

        preserve_existing = bool(existing and existing["status"] == "completed")
        # A rebuild leaves the completed registry row untouched until the
        # staged collection is verified and activated.
        if not preserve_existing:
            self._upsert_registry(
                file_hash, file_name, file_path, file_size, file_ext,
                status="processing",
            )

        try:
            # ===== Step 1: 解析 =====
            if progress_callback:
                progress_callback("解析文件", 0.0)

            file_meta = self._build_file_meta(file_path, file_hash, domain, category)
            file_meta = inherit_registry_metadata(
                file_meta,
                existing,
                explicit_domain=domain,
                explicit_category=category,
            )
            file_ext_lower = file_meta.get("extension", "").lower()

            # PDF + OCR 场景: 渐进入库，避免全部OCR完才嵌入
            if file_ext_lower == ".pdf" and not file_meta.get("is_drawing"):
                ocr_cfg = self.config.get("ocr", {})
                if ocr_cfg.get("enabled"):
                    progressive_result = self._process_pdf_progressive(
                        file_path, file_meta, file_hash, file_name,
                        file_size, file_ext, result, progress_callback,
                        preserve_existing=preserve_existing,
                    )
                    if progressive_result is not None:
                        return progressive_result  # 已处理完成（成功或失败）

            # 通用路径: 解析 → 嵌入 → 入库
            chunks = self._parse_file(file_path, file_meta)
            result.chunks_created = len(chunks)
            result.chars_extracted = sum(c.char_count for c in chunks)
            result.domain = file_meta.get("domain", "")
            result.category = file_meta.get("category", "")
            result.doc_number = file_meta.get("doc_number", "")
            t_parse = (time.time() - t_start) * 1000
            result.parse_time_ms = t_parse

            if not chunks:
                result.status = FileStatus.FAILED
                result.error_message = "解析后无有效文本内容"
                if not preserve_existing:
                    self._upsert_registry(
                        file_hash, file_name, file_path, file_size, file_ext,
                        status="failed", error=result.error_message,
                    )
                return result

            if progress_callback:
                progress_callback("解析文件", 1.0)

            # ===== Step 2: 嵌入 =====
            if progress_callback:
                progress_callback("生成嵌入向量", 0.0)

            result = self._embed_and_insert(
                chunks, result, file_hash, file_name, file_path,
                file_size, file_ext, t_start, progress_callback
            )
            return result

        except MemoryError:
            result.status = FileStatus.FAILED
            result.error_message = "内存不足: 文档过大导致chunker递归溢出, 请尝试拆分文件后重新入库"
            if not preserve_existing:
                self._upsert_registry(file_hash, file_name, file_path, file_size, file_ext,
                                      status="failed", error=result.error_message)
        except Exception as e:
            result.status = FileStatus.FAILED
            result.error_message = str(e) or type(e).__name__
            if not result.error_message.strip():
                result.error_message = f"未知错误: {type(e).__name__}"
            result.error_message = result.error_message[:500]
            if not preserve_existing:
                self._upsert_registry(file_hash, file_name, file_path, file_size, file_ext,
                                      status="failed", error=result.error_message)

        return result

    def process_batch(self, file_paths: List[str],
                      domain: str = None,
                      category: str = None,
                      progress_callback: Callable[[str, float], None] = None,
                      ) -> BatchResult:
        """
        批量处理多个文件

        Args:
            file_paths: 文件路径列表
            domain: 统一指定域
            category: 统一指定类目
            progress_callback: 总体进度回调

        Returns:
            BatchResult
        """
        t_start = time.time()
        results = []
        success = 0
        failed = 0

        for i, fp in enumerate(file_paths):
            if progress_callback:
                progress_callback(f"处理中 ({i + 1}/{len(file_paths)})",
                                  i / len(file_paths))

            result = self.process(fp, domain=domain, category=category)
            results.append(result)
            if result.status == FileStatus.COMPLETED:
                success += 1
            else:
                failed += 1

        if progress_callback:
            progress_callback("完成", 1.0)

        return BatchResult(
            total=len(file_paths),
            success=success,
            failed=failed,
            results=results,
            total_time_ms=(time.time() - t_start) * 1000,
        )

    def delete(self, identifier: str, remove_file: bool = False) -> bool:
        """
        从向量库中删除文件

        Args:
            identifier: 文件 hash 或文件路径
            remove_file: 是否同时删除物理文件（仅限 uploads 目录下的文件）

        Returns:
            是否成功
        """
        file_hash = self._resolve_hash(identifier)
        if not file_hash:
            return False

        with self._metadata_lock, self._index_lock:
            # 获取注册表信息（用于后续清理物理文件）
            reg = self._get_registry(file_hash)

            # Delete from the active alias *and* retained rollback generations.
            # Keeping an old generation here caused re-uploading identical
            # bytes to resurrect its former OCR edits.
            purge = getattr(self.store, "purge_file_generations", None)
            if callable(purge):
                purge(file_hash)
            else:
                self.store.delete_by_file_hash(file_hash)

            # Layout caches are keyed by content hash, so both the current
            # renderer cache and the published-edit sidecar must go together.
            cache_dir = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache"))
            cache_paths = [
                cache_dir / f"{file_hash}.layout.json",
                cache_dir / f"{file_hash}.layout.edited.json",
                cache_dir / f"{file_hash}.layout.json.tmp",
                cache_dir / f"{file_hash}.layout.edited.json.tmp",
            ]
            # A process interrupted during an edit can leave only its hidden
            # publish/restore artifacts behind. They are not valid preview
            # sources, but removing them makes deletion genuinely hash-clean.
            cache_paths.extend(cache_dir.glob(f".{file_hash}.layout*.pending"))
            cache_paths.extend(cache_dir.glob(f".{file_hash}.layout*.restore"))
            for cache_path in cache_paths:
                cache_path.unlink(missing_ok=True)

            # 清理物理文件（安全策略：仅删除 uploads 目录下的文件）
            if remove_file and reg:
                file_path = reg.get("original_path") or reg.get("stored_path") or ""
                if file_path and os.path.exists(file_path):
                    try:
                        uploads_abs = str(self.uploads_dir.resolve())
                        file_abs = str(Path(file_path).resolve())
                        if file_abs.startswith(uploads_abs):
                            os.remove(file_path)
                    except Exception:
                        pass

            # 注册表标记。重新上传相同内容时会重新走 OCR 和入库流程。
            self._upsert_registry(file_hash, status="deleted")
        return True

    def sync_orphans(self, dry_run: bool = False,
                     check_milvus: bool = False) -> dict:
        """
        扫描注册表中指向不存在物理文件的孤记录，清理向量 + 标记 deleted

        Args:
            dry_run: True 只扫描不清理
            check_milvus: 同时检查 Milvus 中是否存在 SQLite 未记录的 chunk 并清理

        Returns:
            {total_checked, orphan_count, cleaned, errors, orphans (仅 dry_run)}
        """
        conn = self._get_db_connection()
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute(
            "SELECT * FROM file_registry WHERE status NOT IN ('deleted', 'pending', 'processing')"
        )
        rows = [dict(row) for row in cursor.fetchall()]
        conn.close()

        orphans = []
        for row in rows:
            file_path = row.get("original_path") or row.get("stored_path") or ""
            if not file_path or not os.path.exists(file_path):
                orphans.append(row)

        # Milvus 一致性检查: 向量库中有 chunk 但 SQLite 无 completed 记录 → 孤 chunk
        milvus_orphan_hashes = []
        sqlite_orphan_hashes = []
        if check_milvus:
            try:
                self.store.client.load_collection(self.store.COLLECTION_NAME)
                raw = self.store.client.query(
                    collection_name=self.store.COLLECTION_NAME,
                    filter="",
                    output_fields=["chunk_id"],
                    limit=30000,
                )
                milvus_hashes = set()
                for r in raw:
                    # chunk_id 格式: "{file_hash}_{chunk_index}"
                    h = r.get("chunk_id", "").rsplit("_", 1)[0]
                    if len(h) == 64:  # SHA256
                        milvus_hashes.add(h)

                conn = self._get_db_connection()
                conn.row_factory = sqlite3.Row
                cursor = conn.cursor()
                cursor.execute(
                    "SELECT file_hash FROM file_registry WHERE status IN ('completed','deleted')"
                )
                sqlite_hashes = {row["file_hash"] for row in cursor.fetchall()}
                conn.close()

                milvus_orphan_hashes = list(milvus_hashes - sqlite_hashes)
                # 反向检查: SQLite completed 但 Milvus 中无 chunk
                for h in [r["file_hash"] for r in rows if r.get("status") == "completed"]:
                    if h not in milvus_hashes:
                        sqlite_orphan_hashes.append(h)
            except Exception as e:
                print(f"[sync_orphans] Milvus query failed: {e}")

        if dry_run:
            result = {
                "total_checked": len(rows),
                "orphan_count": len(orphans),
                "orphans": [
                    {
                        "file_name": o.get("file_name", ""),
                        "file_hash": o.get("file_hash", ""),
                        "original_path": o.get("original_path", "") or o.get("stored_path", ""),
                        "chunks_count": o.get("chunks_count", 0),
                        "domain": o.get("domain", ""),
                        "category": o.get("category", ""),
                    }
                    for o in orphans
                ],
                "cleaned": 0,
            }
            if check_milvus:
                result["milvus_orphan_count"] = len(milvus_orphan_hashes)
                result["milvus_orphan_hashes"] = milvus_orphan_hashes
                result["sqlite_orphan_count"] = len(sqlite_orphan_hashes)
                result["sqlite_orphan_hashes"] = sqlite_orphan_hashes
            return result

        cleaned = 0
        errors = []
        for row in orphans:
            file_hash = row.get("file_hash")
            file_name = row.get("file_name", "")
            try:
                self.store.delete_by_file_hash(file_hash)
                self._upsert_registry(file_hash, status="deleted")
                cleaned += 1
            except Exception as e:
                errors.append({
                    "file_name": file_name,
                    "file_hash": file_hash,
                    "error": str(e),
                })

        # Milvus 孤 chunk 清理
        milvus_cleaned = 0
        if check_milvus and milvus_orphan_hashes:
            for h in milvus_orphan_hashes:
                try:
                    self.store.delete_by_file_hash(h)
                    milvus_cleaned += 1
                except Exception as e:
                    errors.append({
                        "file_hash": h,
                        "error": f"Milvus orphan delete: {e}",
                    })

        # SQLite 孤记录清理: 注册表 completed 但 Milvus 无 chunk → 标记 pending
        sqlite_cleaned = 0
        if check_milvus and sqlite_orphan_hashes:
            for h in sqlite_orphan_hashes:
                try:
                    self._upsert_registry(h, status="pending",
                                          error="Milvus中无对应chunk, 需重新入库")
                    sqlite_cleaned += 1
                except Exception as e:
                    errors.append({
                        "file_hash": h,
                        "error": f"SQLite orphan update: {e}",
                    })

        return {
            "total_checked": len(rows),
            "orphan_count": len(orphans),
            "cleaned": cleaned,
            "errors": errors,
            "milvus_orphan_cleaned": milvus_cleaned if check_milvus else 0,
            "sqlite_orphan_cleaned": sqlite_cleaned if check_milvus else 0,
        }

    def reindex(self, identifier: str,
                progress_callback: Callable = None) -> ProcessResult:
        """
        通过影子集合重建文件索引，验证后原子切换活动别名。

        Args:
            identifier: 文件 hash 或文件路径
        """
        file_hash = self._resolve_hash(identifier)
        if not file_hash:
            return ProcessResult(
                file_path=identifier, file_hash="", file_name="",
                status=FileStatus.FAILED, error_message="文件未在注册表中找到"
            )

        reg = self._get_registry(file_hash)
        if not reg:
            return ProcessResult(
                file_path=identifier, file_hash=file_hash, file_name="",
                status=FileStatus.FAILED, error_message="文件注册信息丢失"
            )

        file_path = reg.get("original_path") or reg.get("stored_path")
        if not file_path or not os.path.exists(file_path):
            return ProcessResult(
                file_path=file_path or "", file_hash=file_hash,
                file_name=reg.get("file_name", ""),
                status=FileStatus.FAILED, error_message="原始文件不存在，无法重建索引"
            )

        return self.process(
            file_path,
            domain=reg.get("domain") or None,
            category=reg.get("category") or None,
            progress_callback=progress_callback,
            force_reindex=True,
        )

    def save_layout_edits(self, identifier: str, base_revision: str,
                          edits: list) -> dict:
        """Publish preview edits to layout cache and Milvus as one operation.

        The original PDF/image remains immutable.  A staged Milvus generation
        is validated before publication; cache files and the active alias are
        compensated back to their previous versions if any publish step fails.
        """
        file_hash = self._resolve_hash(identifier)
        if not file_hash:
            raise LayoutEditError("文件未找到")
        registry = self._get_registry(file_hash)
        if not registry or registry.get("status") != "completed":
            raise LayoutEditError("只有已完成入库的文件可以编辑")

        cache_dir = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache"))
        layout_path = cache_dir / f"{file_hash}.layout.json"
        edited_path = cache_dir / f"{file_hash}.layout.edited.json"
        if not layout_path.exists():
            raise LayoutEditError("当前文件没有可编辑的版面缓存")

        t_start = time.time()
        generation = None
        cache_published = False
        edited_published = False
        layout_before = None
        edited_before = None
        layout_pending = layout_path.with_name(
            f".{layout_path.name}.{threading.get_ident()}.pending"
        )
        edited_pending = edited_path.with_name(
            f".{edited_path.name}.{threading.get_ident()}.pending"
        )

        with self._metadata_lock, self._index_lock:
            try:
                layout_before = layout_path.read_bytes()
                edited_before = edited_path.read_bytes() if edited_path.exists() else None
                current_layout = json.loads(layout_before.decode("utf-8"))
                current_revision = layout_revision(current_layout)
                if not base_revision or base_revision != current_revision:
                    raise LayoutRevisionConflict(
                        "文档已被其他操作更新，请刷新预览后重新编辑"
                    )

                next_layout, audit_rows = apply_layout_edits(current_layout, edits)
                next_revision = layout_revision(next_layout)
                if next_revision == current_revision:
                    raise LayoutEditError("修改内容与当前版本相同")

                file_path = registry.get("original_path") or registry.get("stored_path") or ""
                file_meta = self._build_file_meta(
                    file_path or registry.get("file_name") or "",
                    file_hash,
                    registry.get("domain") or None,
                    registry.get("category") or None,
                )
                file_meta = inherit_registry_metadata(file_meta, registry)
                page_texts = layout_page_texts(next_layout)
                total_pages = max((page for page, _ in page_texts), default=1)
                chunks = []
                for page_num, text in page_texts:
                    if text.strip():
                        chunks.extend(self.chunker.chunk_page_text(
                            text, page_num, total_pages, file_meta
                        ))

                if chunks and self.embedder is None:
                    self._wait_for_gpu_slot("文档编辑向量同步")
                    self.embedder = Embedder(self.config_path)

                generation = self.store.begin_file_generation(file_hash)
                embed_total_ms = 0.0
                for offset in range(0, len(chunks), 20):
                    batch = chunks[offset:offset + 20]
                    embedding_texts = [create_text_for_embedding(chunk) for chunk in batch]
                    t_embed = time.time()
                    encoded = self.embedder.encode(embedding_texts, show_progress=False)
                    embed_total_ms += (time.time() - t_embed) * 1000
                    generation.insert(
                        chunks=batch,
                        dense_vectors=encoded.dense_vectors,
                        sparse_vectors=encoded.sparse_vectors,
                        embedding_texts=embedding_texts,
                    )
                generation.validate(len(chunks), expected_chunks=chunks)

                serialized = json.dumps(next_layout, ensure_ascii=False).encode("utf-8")
                cache_dir.mkdir(parents=True, exist_ok=True)
                layout_pending.write_bytes(serialized)
                edited_pending.write_bytes(serialized)

                generation.activate()
                layout_pending.replace(layout_path)
                cache_published = True
                edited_pending.replace(edited_path)
                edited_published = True

                chars_count = sum(chunk.char_count for chunk in chunks)
                self._upsert_registry(
                    file_hash,
                    registry.get("file_name"),
                    registry.get("original_path"),
                    registry.get("file_size"),
                    registry.get("file_type"),
                    status="completed",
                    chunks_count=len(chunks),
                    chars_count=chars_count,
                    domain=registry.get("domain") or "",
                    category=registry.get("category") or "",
                    doc_number=registry.get("doc_number") or "",
                    parse_time=(time.time() - t_start) * 1000,
                    embed_time=embed_total_ms,
                )
                generation.finalize()
                return {
                    "success": True,
                    "file_hash": file_hash,
                    "revision": next_revision,
                    "changes": audit_rows,
                    "chunks_created": len(chunks),
                    "chars_indexed": chars_count,
                    "visual_assets": count_visual_assets(next_layout),
                    "total_time_ms": (time.time() - t_start) * 1000,
                }
            except BaseException:
                rollback_error = None
                try:
                    if generation is not None:
                        generation.rollback()
                except BaseException as exc:
                    rollback_error = exc
                finally:
                    # Cache restoration must still run if Milvus rollback itself
                    # reports an error; otherwise the two stores could diverge.
                    if cache_published and layout_before is not None:
                        restore = layout_path.with_name(f".{layout_path.name}.restore")
                        restore.write_bytes(layout_before)
                        restore.replace(layout_path)
                    if edited_published:
                        if edited_before is None:
                            edited_path.unlink(missing_ok=True)
                        else:
                            restore = edited_path.with_name(f".{edited_path.name}.restore")
                            restore.write_bytes(edited_before)
                            restore.replace(edited_path)
                if rollback_error is not None:
                    print(f"[EDIT] Milvus rollback failed: {rollback_error}", flush=True)
                raise
            finally:
                layout_pending.unlink(missing_ok=True)
                edited_pending.unlink(missing_ok=True)

    def list_files(self, status: str = None, domain: str = None,
                   limit: int = 100, offset: int = 0,
                   check_existence: bool = True,
                   exclude_deleted: bool = False) -> List[dict]:
        """列出已注册的文件

        Args:
            check_existence: 检查物理文件是否存在，添加 file_exists 字段
            exclude_deleted: 排除已删除文件 (仅在未显式指定 status 时生效)
        """
        conn = self._get_db_connection()
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()

        query = "SELECT * FROM file_registry WHERE 1=1"
        params = []
        if status:
            query += " AND status = ?"
            params.append(status)
        elif exclude_deleted:
            query += " AND status != 'deleted'"
        if domain:
            query += " AND domain = ?"
            params.append(domain)

        query += " ORDER BY updated_at DESC LIMIT ? OFFSET ?"
        params.extend([limit, offset])

        cursor.execute(query, params)
        rows = [dict(row) for row in cursor.fetchall()]
        conn.close()

        # 检测物理文件是否存在
        for row in rows:
            file_path = (row.get("original_path") or row.get("stored_path") or "")
            if file_path and os.path.exists(file_path):
                row["file_exists"] = True
            elif file_path and not os.path.exists(file_path):
                row["file_exists"] = False
            else:
                row["file_exists"] = None  # 无路径信息

        return rows

    def update_file_meta(self, file_hash: str, updates: dict) -> dict:
        """原子地同步注册表与 Milvus chunk 元数据。"""
        filtered = normalize_metadata_updates(updates)
        if not filtered:
            return {"success": False, "error": "无可更新的字段"}
        resolved_hash = self._resolve_hash(file_hash)
        if not resolved_hash:
            return {"success": False, "error": "文件未找到"}

        with self._metadata_lock:
            # Hold SQLite's writer lock across the Milvus operation so another
            # process cannot concurrently reindex or edit the same registry row.
            conn = self._get_db_connection()
            conn.row_factory = sqlite3.Row
            try:
                conn.execute("BEGIN IMMEDIATE")
                registry_row = conn.execute(
                    "SELECT * FROM file_registry WHERE file_hash = ?",
                    (resolved_hash,),
                ).fetchone()
                if not registry_row:
                    conn.rollback()
                    return {"success": False, "error": "文件未找到"}
                registry = dict(registry_row)

                def commit_registry(values):
                    sets = ", ".join(f"{key} = ?" for key in values)
                    params = [*values.values(), resolved_hash]
                    cursor = conn.execute(
                        f"UPDATE file_registry SET {sets}, "
                        "updated_at = datetime('now','localtime') "
                        "WHERE file_hash = ?",
                        params,
                    )
                    if cursor.rowcount != 1:
                        raise RuntimeError("File registry row disappeared during update")
                    conn.commit()

                chunk_count = synchronize_metadata(
                    file_hash=resolved_hash,
                    updates=filtered,
                    expected_chunks=int(registry.get("chunks_count") or 0),
                    store=self.store,
                    commit_registry=commit_registry,
                )
            except BaseException:
                conn.rollback()
                raise
            finally:
                conn.close()
        return {
            "success": True,
            "updated": 1,
            "chunks_updated": chunk_count,
            "file_hash": resolved_hash,
        }

    def get_distinct_subcategories(self, domain: str = None, category: str = None) -> list:
        """返回去重子类目列表，供前端级联选择器"""
        conn = self._get_db_connection()
        cursor = conn.cursor()
        query = "SELECT DISTINCT domain, category FROM file_registry WHERE status = 'completed'"
        params = []
        if domain:
            query += " AND domain = ?"
            params.append(domain)
        if category:
            query += " AND category = ?"
            params.append(category)
        query += " ORDER BY domain, category"
        cursor.execute(query, params)
        rows = [{"domain": r[0], "category": r[1]} for r in cursor.fetchall()]
        conn.close()
        return rows

    def count_files(self, status: str = None, domain: str = None, exclude_deleted: bool = False) -> int:
        """返回匹配条件的文件总数 (不受 limit/offset 限制)"""
        conn = self._get_db_connection()
        cursor = conn.cursor()
        query = "SELECT COUNT(*) FROM file_registry WHERE 1=1"
        params = []
        if status:
            query += " AND status = ?"
            params.append(status)
        elif exclude_deleted:
            query += " AND status != 'deleted'"
        if domain:
            query += " AND domain = ?"
            params.append(domain)
        cursor.execute(query, params)
        count = cursor.fetchone()[0]
        conn.close()
        return count

    def get_summary(self) -> dict:
        """获取索引入库摘要

        统计规则:
          - total_files: 仅统计已入库完成 (completed) 的文件
          - total_chunks: 优先从 Milvus 向量库取真实值，失败回退到注册表
          - by_status: 保留完整状态分布供调试
        """
        conn = self._get_db_connection()
        cursor = conn.cursor()

        cursor.execute("SELECT status, COUNT(*) FROM file_registry GROUP BY status")
        by_status = {row[0]: row[1] for row in cursor.fetchall()}

        cursor.execute("SELECT domain, COUNT(*) FROM file_registry WHERE status='completed' GROUP BY domain")
        by_domain = {row[0]: row[1] for row in cursor.fetchall()}

        cursor.execute("SELECT SUM(chunks_count), SUM(chars_count) FROM file_registry WHERE status='completed'")
        totals = cursor.fetchone()

        conn.close()

        # total_files: 只统计已入库完成 (completed) 的文件
        total_files = by_status.get("completed", 0)

        # total_chunks: 优先从 Milvus 向量库取真实值
        registry_chunks = totals[0] or 0
        total_chunks = registry_chunks
        try:
            store_stats = self.store.get_collection_stats()
            if store_stats.get("exists"):
                total_chunks = store_stats.get("count", 0)
        except Exception:
            pass  # Milvus 不可用时回退到注册表统计

        return {
            "total_files": total_files,
            "by_status": by_status,
            "by_domain": by_domain,
            "total_chunks": total_chunks,
            "total_chars": totals[1] or 0,
        }

    # ===== 内部方法 =====

    def _get_registry(self, file_hash: str) -> Optional[dict]:
        conn = self._get_db_connection()
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM file_registry WHERE file_hash = ?", (file_hash,))
        row = cursor.fetchone()
        conn.close()
        return dict(row) if row else None

    def _upsert_registry(self, file_hash: str, file_name: str = None,
                         original_path: str = None, file_size: int = None,
                         file_type: str = None, status: str = None,
                         chunks_count: int = None, chars_count: int = None,
                         domain: str = None, category: str = None,
                         doc_number: str = None, error: str = None,
                         parse_time: float = None, embed_time: float = None):
        conn = self._get_db_connection()
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()

        # 内联 SELECT，避免 _get_registry 开第二个连接导致 database is locked
        cursor.execute("SELECT * FROM file_registry WHERE file_hash = ?", (file_hash,))
        row = cursor.fetchone()
        existing = dict(row) if row else None

        if existing:
            updates = []
            params = []
            for col, val in [
                ("original_path", original_path), ("file_name", file_name),
                ("file_size", file_size), ("file_type", file_type), ("status", status),
                ("chunks_count", chunks_count), ("chars_count", chars_count),
                ("domain", domain), ("category", category), ("doc_number", doc_number),
                ("error_message", error), ("parse_time_ms", parse_time),
                ("embed_time_ms", embed_time),
            ]:
                if val is not None:
                    updates.append(f"{col} = ?")
                    params.append(val)
            # 成功后清除旧错误信息
            if status == "completed":
                updates.append("error_message = ''")
            if status == "completed" and existing["status"] == "completed":
                updates.append("reindex_count = reindex_count + 1")
            updates.append("updated_at = datetime('now', 'localtime')")
            params.append(file_hash)
            cursor.execute(f"UPDATE file_registry SET {', '.join(updates)} WHERE file_hash = ?", params)
        else:
            cursor.execute("""
                INSERT INTO file_registry (file_hash, original_path, file_name,
                file_size, file_type, status, chunks_count, chars_count,
                domain, category, doc_number, error_message,
                parse_time_ms, embed_time_ms)
                VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)
            """, (file_hash, original_path or "", file_name or "", file_size or 0,
                  file_type or "", status or "pending", chunks_count or 0,
                  chars_count or 0, domain or "", category or "",
                  doc_number or "", error or "", parse_time or 0, embed_time or 0))

        conn.commit()
        conn.close()

    def _process_pdf_progressive(self, file_path, file_meta, file_hash, file_name,
                                  file_size, file_ext, result: ProcessResult,
                                  progress_callback=None,
                                  preserve_existing: bool = False) -> Optional[ProcessResult]:
        """
        PDF 按页入库 — 保留页码, 结构化分块

        流程:
          1. 解析 PDF, 逐页判定 needs_ocr
          2. 每页: 纯文本页用 fitz 文本, 扫描页调用外部 PaddleOCR-VL 结构化识别
          3. 按页分块 (chunk_page_text) → 保留 page_num, chunk_id 无冲突
          4. 嵌入 + 写入 Milvus

        OCR/LLM 均为外部服务, 无需本地显存调度
        """
        t_start = time.time()
        parsed = self.pdf_parser.parse(file_path)
        page_count = parsed.get("page_count", 0)
        needs_ocr = parsed.get("needs_ocr_pages", [])
        cache_dir = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache"))
        edited_layout_path = cache_dir / f"{file_hash}.layout.edited.json"
        edited_page_texts = {}
        if edited_layout_path.exists():
            # A user-edited layout is the published source of truth.  Reindexing
            # must rebuild vectors from it rather than silently restoring OCR
            # text/images from the immutable original PDF.
            try:
                edited_layout = json.loads(edited_layout_path.read_text(encoding="utf-8"))
                edited_page_texts = dict(layout_page_texts(edited_layout))
                cache_dir.mkdir(parents=True, exist_ok=True)
                (cache_dir / f"{file_hash}.layout.json").write_text(
                    json.dumps(edited_layout, ensure_ascii=False), encoding="utf-8"
                )
                needs_ocr = []
                print(
                    f"   [EDIT] 保留已发布编辑版面，共 {len(edited_page_texts)} 页",
                    flush=True,
                )
            except (OSError, ValueError, LayoutEditError) as exc:
                raise RuntimeError(f"已编辑版面缓存损坏，拒绝覆盖: {exc}") from exc
        elif self.config.get("ocr", {}).get("always_pipeline"):
            needs_ocr = list(range(1, page_count + 1))
            print(f"   [OCR] always_pipeline=true，8001 Pipeline 处理全部 {page_count} 页", flush=True)

        # 批量识别所有需要 OCR 的页面 (VL 服务逐页调用)
        ocr_texts = {}
        if needs_ocr:
            if progress_callback:
                progress_callback(f"OCR识别 ({len(needs_ocr)}页)", 0.1)
            print(f"   [OCR] 外部 VL 服务识别 {len(needs_ocr)} 页...")
            try:
                ocr_texts = self.pdf_parser.ocr_pages(
                    file_path, [p - 1 for p in needs_ocr]
                )
                layout_pages = getattr(self.pdf_parser._vl_client, "last_layout_pages", {})
                if layout_pages:
                    cache_dir = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache"))
                    cache_dir.mkdir(parents=True, exist_ok=True)
                    layout_path = cache_dir / f"{file_hash}.layout.json"
                    layout_path.write_text(
                        json.dumps(layout_pages, ensure_ascii=False), encoding="utf-8"
                    )
                    print(f"   [OCR] 已缓存 {len(layout_pages)} 页版面块: {layout_path.name}")
            except Exception as exc:
                # OCR is an enhancement. Retain the parser's existing text for
                # every affected page when the external OCR service fails.
                print(f"   [OCR] 服务异常，回退 PDF 原文本: {exc}")
                ocr_texts = {}
            if progress_callback:
                progress_callback("OCR完成, 分块中...", 0.5)

        # 按页分块: 每页保留 page_num, chunk_id 唯一
        all_chunks = []
        total_chars = 0
        # 垃圾 OCR 检测 (小模型幻觉/重复输出)
        from ingestion.vl_ocr import is_garbage_ocr_text

        for page in parsed.get("pages", []):
            pn = page.get("page_num", 1)
            if edited_page_texts:
                text = edited_page_texts.get(pn, "")
            elif page.get("needs_ocr"):
                text, text_source = choose_page_text(
                    page.get("text", ""),
                    ocr_texts.get(pn - 1, ""),
                    is_garbage_ocr_text,
                )
                if text_source == "pdf_text_fallback":
                    print(f"   [OCR] 第{pn}页 OCR 无效，已保留 PDF 原文本")
            else:
                text = page.get("text", "")
            if not (text or "").strip():
                continue
            page_chunks = self.chunker.chunk_page_text(
                text, pn, page_count, file_meta
            )
            all_chunks.extend(page_chunks)
            total_chars += sum(c.char_count for c in page_chunks)

        # ponytail: parsed / ocr_texts 已消费, 释放
        del parsed, ocr_texts

        if not all_chunks:
            result.status = FileStatus.FAILED
            result.error_message = "解析后无有效文本内容 (含OCR)"
            if not preserve_existing:
                self._upsert_registry(
                    file_hash, file_name, file_path, file_size,
                    file_ext, status="failed", error=result.error_message,
                )
            return result

        result.chunks_created = len(all_chunks)
        result.chars_extracted = total_chars
        result.domain = file_meta.get("domain", "")
        result.category = file_meta.get("category", "")
        result.doc_number = file_meta.get("doc_number", "")
        result.parse_time_ms = (time.time() - t_start) * 1000

        # 统一走影子集合嵌入路径，活动集合在完整验证前保持不变。
        return self._embed_and_insert(all_chunks, result, file_hash, file_name,
                                      file_path, file_size, file_ext,
                                      t_start, progress_callback)

    def _embed_and_insert(self, chunks, result: ProcessResult,
                          file_hash, file_name, file_path, file_size, file_ext,
                          t_start, progress_callback=None) -> ProcessResult:
        """分批嵌入 + 逐批写入 — ponytail: 时间换空间, 控制内存峰值"""
        if self.embedder is None:
            self._wait_for_gpu_slot("BGE-M3 嵌入模型加载")
            self.embedder = Embedder(self.config_path)

        # ponytail: 每批20个chunk, 内存峰值 = 20组文本+向量, 而非全部
        BATCH = 20
        total = len(chunks)
        embed_total_ms = 0.0
        generation = None
        with self._index_lock:
            try:
                generation = self.store.begin_file_generation(file_hash)
                for i in range(0, total, BATCH):
                    batch = chunks[i:i + BATCH]
                    batch_texts = [create_text_for_embedding(c) for c in batch]

                    if progress_callback:
                        end = min(i + BATCH, total)
                        progress_callback(
                            f"嵌入 ({i + 1}-{end}/{total})",
                            i / total * 0.9,
                        )

                    t_embed = time.time()
                    emb_result = self.embedder.encode(
                        batch_texts, show_progress=False
                    )
                    embed_total_ms += (time.time() - t_embed) * 1000
                    generation.insert(
                        chunks=batch,
                        dense_vectors=emb_result.dense_vectors,
                        sparse_vectors=emb_result.sparse_vectors,
                        embedding_texts=batch_texts,
                    )
                    del emb_result, batch_texts

                generation.validate(total, expected_chunks=chunks)
                generation.activate()

                result.embed_time_ms = embed_total_ms
                result.status = FileStatus.COMPLETED
                result.total_time_ms = (time.time() - t_start) * 1000

                # The registry becomes completed only after the alias switch.
                # If this commit fails, rollback switches the alias back.
                self._upsert_registry(
                    file_hash, file_name, file_path, file_size, file_ext,
                    status="completed",
                    chunks_count=len(chunks),
                    chars_count=result.chars_extracted,
                    domain=result.domain,
                    category=result.category,
                    doc_number=result.doc_number,
                    parse_time=result.parse_time_ms,
                    embed_time=result.embed_time_ms,
                )
                generation.finalize()
            except BaseException:
                if generation is not None:
                    generation.rollback()
                raise
        return result

    def _wait_for_gpu_slot(self, task_name: str = "入库", min_free_mb: int = 2500):
        """GPU 显存背压 — 等待足够显存后再执行操作"""
        try:
            sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
            from utils.gpu_monitor import get_gpu_monitor
            monitor = get_gpu_monitor(min_free_vram_mb=min_free_mb)
            ok = monitor.wait_for_vram(min_free_mb=min_free_mb)
            if not ok:
                print(f"   [WARN] [{task_name}] 显存等待超时，强制执行")
            return ok
        except ImportError:
            pass
        except Exception:
            pass
        return True

    def _resolve_hash(self, identifier: str) -> Optional[str]:
        """从文件路径或 hash 解析为 hash"""
        if len(identifier) == 64 and all(c in "0123456789abcdef" for c in identifier):
            return identifier
        if os.path.exists(identifier):
            return self.compute_hash(identifier)
        # 从注册表查找
        conn = self._get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT file_hash FROM file_registry WHERE original_path = ? OR file_name = ?",
            (identifier, Path(identifier).name)
        )
        row = cursor.fetchone()
        conn.close()
        return row[0] if row else None

    def _build_file_meta(self, file_path: str, file_hash: str,
                         domain: str = None, category: str = None) -> dict:
        """构建元数据字典"""
        from ingestion.file_walker import FileWalker

        fp = Path(file_path)
        # 如果文件在知识库目录中，用 FileWalker 提取元数据
        kb_path = self.config["paths"]["knowledge_base"]
        walker = FileWalker(self.config_path)

        try:
            rel_path = str(fp.relative_to(kb_path))
        except ValueError:
            rel_path = fp.name

        path_meta = walker.extract_path_metadata(rel_path)
        filename_meta = walker.extract_filename_metadata(fp.name, str(fp))

        meta = {
            "file_hash": file_hash,
            "full_path": str(fp),
            "relative_path": rel_path,
            "file_name": fp.name,
            "extension": fp.suffix.lower(),
            "size_bytes": fp.stat().st_size if fp.exists() else 0,
            "domain": domain or path_meta.get("domain") or "",
            "category": category or path_meta.get("category") or "",
            "subcategory": path_meta.get("subcategory") or "",
            "doc_number": filename_meta.get("doc_number") or "",
            "publish_level": filename_meta.get("publish_level") or "",
            "voltage_level": filename_meta.get("voltage_level") or "",
            "discipline": filename_meta.get("discipline") or "",
            "equipment_type": filename_meta.get("equipment_type") or "",
            "year": filename_meta.get("year") or 0,
            "region": filename_meta.get("region") or "全国",
            "drawing_code": filename_meta.get("drawing_code") or "",
        }

        # 判断是否图纸
        is_dwg = meta["extension"] in (".dwg", ".dxf")
        meta["is_drawing"] = 1 if is_dwg else 0
        meta["is_archive"] = 1 if meta["extension"] in (".zip", ".rar", ".7z") else 0
        meta["format_group"] = "drawing" if is_dwg else "document"

        return meta

    def _parse_file(self, file_path: str, file_meta: dict) -> List[Chunk]:
        """解析文件并生成 chunks"""
        ext = file_meta["extension"].lower()
        is_drawing = file_meta.get("is_drawing", 0)

        # PDF
        if ext == ".pdf":
            if is_drawing:
                text = self.pdf_parser.parse_single_page_pdf(file_path) or ""
                return self.chunker.chunk_drawing(text, file_meta)
            else:
                # ponytail: OCR 路径已在 _process_pdf_progressive 处理,
                # 此分支仅当 OCR 禁用时走到
                parsed = self.pdf_parser.parse(file_path)
                return self.chunker.chunk_pdf_document(parsed, file_meta)

        # 图片：上传接口允许 PNG/JPEG，但此前只在 PDF 页面路径中调用
        # VL-OCR，导致独立图片落入“其他”分支并被判定为无有效文本。
        # 复用同一 OCR 客户端（当前配置为 8001 PaddleOCR pipeline），
        # 保持与 PDF OCR 相同的 Markdown/表格识别结果。
        elif ext in (".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".gif"):
            cache_dir = Path(self.config["paths"].get("parsed_cache", "data/parsed_cache"))
            edited_path = cache_dir / f"{file_meta.get('file_hash')}.layout.edited.json"
            if edited_path.exists():
                try:
                    layout_blocks = json.loads(edited_path.read_text(encoding="utf-8"))
                    page_text = layout_page_texts(layout_blocks)
                    text = "\n\n".join(value for _, value in page_text if value.strip())
                    (cache_dir / f"{file_meta.get('file_hash')}.layout.json").write_text(
                        json.dumps(layout_blocks, ensure_ascii=False), encoding="utf-8"
                    )
                    print("   [EDIT] 图片重新入库保留已发布编辑版面", flush=True)
                    return self._chunk_text_safe(text, file_meta) if text.strip() else []
                except (OSError, ValueError, LayoutEditError) as exc:
                    raise RuntimeError(f"已编辑版面缓存损坏，拒绝覆盖: {exc}") from exc
            print(
                f"[OCR-TRACE] image branch file={file_path} ext={ext} "
                f"ocr_enabled={self.config.get('ocr', {}).get('enabled')} "
                f"protocol={self.config.get('ocr', {}).get('vl', {}).get('protocol')} "
                f"base_url={self.config.get('ocr', {}).get('vl', {}).get('base_url')}",
                flush=True,
            )
            text, layout_blocks = self.pdf_parser.ocr_page_with_layout(file_path)
            if layout_blocks:
                cache_dir.mkdir(parents=True, exist_ok=True)
                layout_path = cache_dir / f"{file_meta.get('file_hash')}.layout.json"
                layout_path.write_text(json.dumps(layout_blocks, ensure_ascii=False), encoding="utf-8")
                print(f"   [OCR] 已缓存版面块 {len(layout_blocks)} 个: {layout_path.name}")
            if not text.strip():
                return []
            return self._chunk_text_safe(text, file_meta)

        # DOC (old binary Word format)
        elif ext == ".doc":
            text = self._parse_doc_file(file_path)
            if text:
                return self.chunker.chunk_text_document(text, file_meta)
            return []

        # DOCX
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(file_path)

                # 提取段落文本
                para_text = "\n".join([p.text for p in doc.paragraphs])

                # 提取表格内容（电力规范docx大量信息在表格中）
                table_parts = []
                for ti, table in enumerate(doc.tables):
                    rows_text = []
                    for row in table.rows:
                        cells = []
                        for cell in row.cells:
                            cell_text = " ".join(p.text for p in cell.paragraphs if p.text.strip())
                            if cell_text:
                                cells.append(cell_text)
                        if cells:
                            rows_text.append(" | ".join(cells))
                    if rows_text:
                        table_parts.append(
                            f"[表格{ti + 1}]\n" + "\n".join(rows_text)
                        )

                text = para_text
                if table_parts:
                    text += "\n\n" + "\n\n".join(table_parts)

                if not text.strip():
                    return []

                return self.chunker.chunk_text_document(text, file_meta)

            except ImportError:
                print(f"   [warn] python-docx 未安装，无法解析 .docx: {file_path}")
                return []
            except Exception as e:
                print(f"   [warn] .docx 解析失败: {Path(file_path).name}: {e}")
                return []

        # TXT / MD
        elif ext in (".txt", ".md"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return self.chunker.chunk_text_document(text, file_meta)

        # OFD
        elif ext == ".ofd":
            text = self._parse_ofd_file(file_path)
            if not text or not text.strip():
                return []
            return self._chunk_text_safe(text, file_meta)

        # XLSX (Office Open XML)
        elif ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                texts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = [f"[Sheet: {sheet_name}]"]
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join([str(c) if c is not None else "" for c in row])
                        if row_text.strip(" |"):
                            sheet_text.append(row_text)
                    texts.append("\n".join(sheet_text))
                text = "\n\n".join(texts)
                return self.chunker.chunk_text_document(text, file_meta)
            except Exception as e:
                print(f"[warn] XLSX解析失败 ({os.path.basename(file_path)}): {e}")
                return []

        # XLS (OLE2 binary format — openpyxl can't handle this)
        elif ext == ".xls":
            try:
                import xlrd
                wb = xlrd.open_workbook(file_path)
                texts = []
                for sheet_name in wb.sheet_names():
                    ws = wb.sheet_by_name(sheet_name)
                    sheet_text = [f"[Sheet: {sheet_name}]"]
                    for row_idx in range(ws.nrows):
                        row_values = ws.row_values(row_idx)
                        row_text = " | ".join([str(c) if c != "" else "" for c in row_values])
                        if row_text.strip(" |"):
                            sheet_text.append(row_text)
                    texts.append("\n".join(sheet_text))
                text = "\n\n".join(texts)
                return self.chunker.chunk_text_document(text, file_meta)
            except Exception as e:
                print(f"[warn] XLS解析失败 ({os.path.basename(file_path)}): {e}")
                return []

        # WPS (金山 WPS 文字文档)
        # WPS 有两代格式: 旧版 OLE2 二进制容器, 新版 ZIP+XML 容器
        elif ext == ".wps":
            text = self._parse_wps_file(file_path)
            if text:
                return self.chunker.chunk_text_document(text, file_meta)
            return []

        # PPTX (python-pptx 提取幻灯片文本)
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for slide in prs.slides:
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    texts.append(t)
                    if texts:
                        slides_text.append("\n".join(texts))
                text = "\n\n".join(slides_text)
                if text.strip():
                    return self.chunker.chunk_text_document(text, file_meta)
                return []
            except ImportError:
                print("   [WARN] python-pptx 未安装，无法解析 .pptx")
                return []
            except Exception as e:
                print(f"   [WARN] PPTX 解析失败: {e}")
                return []

        # PPT (旧版) — 尝试用 Windows COM 或 LibreOffice
        elif ext == ".ppt":
            text = self._parse_ppt_file(file_path)
            if text:
                return self.chunker.chunk_text_document(text, file_meta)
            return []

        # DWG (跳过，需专用解析器)
        elif ext in (".dwg", ".dxf"):
            return []

        # 其他
        else:
            return []

    def _parse_doc_file(self, file_path: str) -> str:
        """
        解析旧版 .doc 文件 (OLE2 复合文档格式)
        按优先级尝试多种后端:
          1. win32com (Windows MS Word COM 自动化, 最可靠)
          2. LibreOffice headless 转换
          3. olefile 原始文本提取
          4. antiword (Linux)
          5. docx2txt / python-docx (仅对伪装的 .docx 有效)
        """
        # 方案1: Windows COM (MS Word 安装时最可靠)
        text = self._parse_doc_via_win32(file_path)
        if text and text.strip():
            return text

        # 方案2: LibreOffice headless 转换
        text = self._parse_doc_via_libreoffice(file_path)
        if text and text.strip():
            return text

        # 方案3: olefile 原始提取
        text = self._parse_doc_via_olefile(file_path)
        if text and text.strip():
            return text

        # 方案4: antiword (Linux)
        text = self._parse_doc_via_antiword(file_path)
        if text and text.strip():
            return text

        # 方案5: docx2txt / python-docx (某些 .doc 实际是 .docx 改名)
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            if text and text.strip():
                print(f"   [doc] docx2txt 解析成功: {len(text)} 字符")
                return text
        except Exception:
            pass

        try:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            if text and text.strip():
                print(f"   [doc] python-docx 解析成功: {len(text)} 字符")
                return text
        except Exception:
            pass

        # 方案6: raw binary UTF-16LE 兜底 (纯二进制提取, 零依赖)
        text = self._parse_doc_via_raw_binary(file_path)
        if text and text.strip():
            print(f"   [doc] raw binary 解析成功: {len(text)} 字符")
            return text

        print(f"   [warn] 所有 .doc 解析方案均失败: {os.path.basename(file_path)}")
        print(f"   [tip] 建议方案: (1) pip install pywin32 启用 Word COM 解析")
        print(f"         或 (2) 安装 LibreOffice")
        print(f"         或 (3) 用 Word 打开后另存为 .docx 格式")
        return ""

    def _parse_doc_via_raw_binary(self, file_path: str) -> str:
        """
        OLE2 .doc 纯二进制兜底提取 — 零外部依赖

        直接以 UTF-16LE 解码文件二进制内容，用正则提取中文文本段。
        适用于 OLE2 容器中文档内容以 Unicode 存储的常见情况。
        对复杂排版（表格/图片/公式）可能丢失部分内容，但绝大多数
        .doc 文件能够提取到 80%+ 的正文文本。
        """
        import re
        try:
            with open(file_path, 'rb') as f:
                data = f.read()
            # OLE2 .doc Word文档内文本以 UTF-16LE 编码存储
            text = data.decode('utf-16-le', errors='ignore')
            # 提取有意义的文本段: 4个字符以上, 中文占比 > 10%
            chunks = re.findall(
                r"[一-鿿a-zA-Z0-9\s.,;:!?()（）、。，；：！？''【】《》——…/%+\-]{4,}",
                text
            )
            # 过滤乱码: 中文占比太低的多半是字体名/样式名等垃圾
            clean = []
            for c in chunks:
                cn_count = sum(1 for ch in c if '一' <= ch <= '鿿')
                ratio = cn_count / max(len(c), 1)
                if ratio > 0.05 and len(c) > 6:
                    clean.append(c)
            return '\n'.join(clean)
        except Exception:
            return ""

    def _parse_doc_via_win32(self, file_path: str) -> str:
        """
        通过 Windows COM 调用 Microsoft Word 提取文本
        这是 Windows 上解析 .doc 最可靠的方式
        """
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            return ""

        abs_path = os.path.abspath(file_path)
        word = None
        doc = None
        try:
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0

            # 打开文档
            doc = word.Documents.Open(abs_path, ReadOnly=True, Visible=False)

            # 提取所有文本
            text = doc.Content.Text

            # 关闭文档
            doc.Close(SaveChanges=False)

            if text and text.strip():
                print(f"   [doc] Word COM 解析成功: {len(text)} 字符")
                return text

        except Exception as e:
            print(f"   [doc] Word COM 失败: {e}")
            # 确保即使出错也尝试关闭文档
            if doc is not None:
                try:
                    doc.Close(SaveChanges=False)
                except Exception:
                    pass
        finally:
            if word is not None:
                try:
                    word.Quit()
                except Exception:
                    pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

        return ""

    def _parse_doc_via_libreoffice(self, file_path: str) -> str:
        """通过 LibreOffice headless 将 .doc 转为文本"""
        import subprocess
        import tempfile

        # 查找 LibreOffice 路径
        lo_paths = [
            "libreoffice", "soffice",
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
            "/usr/bin/libreoffice", "/usr/bin/soffice",
        ]

        lo_exe = None
        for p in lo_paths:
            try:
                subprocess.run([p, "--version"], capture_output=True, timeout=5)
                lo_exe = p
                break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue

        if not lo_exe:
            return ""

        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                cmd = [
                    lo_exe, "--headless", "--convert-to", "txt:Text",
                    "--outdir", tmpdir, file_path,
                ]
                result = subprocess.run(cmd, capture_output=True, timeout=60)
                if result.returncode == 0:
                    # 查找生成的 txt 文件
                    for f in os.listdir(tmpdir):
                        if f.endswith(".txt"):
                            txt_path = os.path.join(tmpdir, f)
                            with open(txt_path, "r", encoding="utf-8", errors="ignore") as fp:
                                text = fp.read()
                            if text.strip():
                                print(f"   [doc] LibreOffice 解析成功: {len(text)} 字符")
                                return text
        except Exception:
            pass

        return ""

    def _parse_doc_via_olefile(self, file_path: str) -> str:
        """通过 olefile 从 OLE2 容器中提取原始文本"""
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)

            # 尝试读取 WordDocument 流中的文本
            text_parts = []

            # 读取主文本流
            if ole.exists("WordDocument"):
                data = ole.openstream("WordDocument").read()
                # 尝试提取可读文本 (UTF-16 LE 编码的文本片段)
                try:
                    decoded = data.decode("utf-16-le", errors="ignore")
                    # 过滤控制字符，保留可读内容
                    import re
                    readable = re.findall(r'[一-鿿　-〿＀-￯a-zA-Z0-9\s.,;:!?()（）、。，；：！？""''【】《》/-]+', decoded)
                    if readable:
                        text_parts.extend(readable)
                except Exception:
                    pass

            # 尝试 1Table 或 0Table 流
            for stream_name in ole.listdir():
                stream_path = "/".join(stream_name) if isinstance(stream_name, list) else stream_name
                if "Table" in stream_path or "Text" in stream_path:
                    try:
                        data = ole.openstream(stream_path).read()
                        decoded = data.decode("utf-16-le", errors="ignore")
                        import re
                        readable = re.findall(r'[一-鿿]+', decoded)
                        if readable:
                            text_parts.extend(readable)
                    except Exception:
                        pass

            ole.close()

            if text_parts:
                text = " ".join(text_parts)
                print(f"   [doc] olefile 解析成功: {len(text)} 字符 (可能不完整)")
                return text
        except ImportError:
            pass
        except Exception:
            pass

        return ""

    def _parse_doc_via_antiword(self, file_path: str) -> str:
        """通过 antiword 解析 (Linux)"""
        import subprocess
        try:
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True, timeout=30,
            )
            if result.returncode == 0:
                text = result.stdout.decode("utf-8", errors="ignore")
                if text.strip():
                    print(f"   [doc] antiword 解析成功: {len(text)} 字符")
                    return text
        except FileNotFoundError:
            pass
        except Exception:
            pass
        return ""

    # ===== WPS 文件解析 (.wps) =====

    def _parse_wps_file(self, file_path: str) -> str:
        """
        解析 .wps 文件 (金山 WPS 文字文档)

        WPS 有两代格式:
          - 新版 WPS (2010+): ZIP + XML 容器, 内部结构与 .docx 相同
          - 旧版 WPS (2005 及以前): OLE2 二进制容器, 类似旧 .doc

        按优先级尝试多种后端:
          1. python-docx (新版 .wps 本质是 .docx)
          2. zipfile 直接解压提取 XML 文本 (新版备选)
          3. LibreOffice headless 转换 (旧版+新版通用)
          4. olefile 原始提取 (旧版 OLE2)
          5. win32com WPS Office COM 自动化 (最可靠, 需安装 WPS Office)
          6. win32com MS Word COM 自动化 (MS Word 可能能打开)
        """
        abs_path = os.path.abspath(file_path)
        fname = os.path.basename(file_path)

        # 方案1: python-docx (新版 .wps = .docx)
        text = self._parse_wps_via_docx(file_path)
        if text and text.strip():
            print(f"   [wps] python-docx 解析成功: {len(text)} 字符")
            return text

        # 方案2: zipfile 直接提取 XML (新版 .wps)
        text = self._parse_wps_via_zip(file_path)
        if text and text.strip():
            print(f"   [wps] zipfile 解析成功: {len(text)} 字符")
            return text

        # 方案3: LibreOffice headless 转换 (通用)
        text = self._parse_doc_via_libreoffice(file_path)
        if text and text.strip():
            print(f"   [wps] LibreOffice 解析成功: {len(text)} 字符")
            return text

        # 方案4: olefile 原始提取 (旧版 .wps OLE2)
        text = self._parse_doc_via_olefile(file_path)
        if text and text.strip():
            print(f"   [wps] olefile 解析成功: {len(text)} 字符 (可能不完整)")
            return text

        # 方案5: WPS Office COM 自动化 (需安装 WPS Office)
        text = self._parse_wps_via_wps_com(file_path)
        if text and text.strip():
            print(f"   [wps] WPS COM 解析成功: {len(text)} 字符")
            return text

        # 方案6: MS Word COM (可能兼容某些 .wps)
        text = self._parse_doc_via_win32(file_path)
        if text and text.strip():
            print(f"   [wps] Word COM 解析成功: {len(text)} 字符")
            return text

        print(f"   [warn] 所有 .wps 解析方案均失败: {fname}")
        print(f"   [tip] 建议方案: (1) 安装 WPS Office 启用 COM 解析")
        print(f"         或 (2) 安装 LibreOffice")
        print(f"         或 (3) 用 WPS 打开后另存为 .docx 格式")
        return ""

    def _parse_wps_via_docx(self, file_path: str) -> str:
        """
        尝试用 python-docx 解析 .wps。
        新版 WPS 文件本质上是 ZIP + XML，与 .docx 结构相同。
        """
        try:
            import docx
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            if paragraphs:
                return "\n".join(paragraphs)
        except Exception:
            pass
        return ""

    def _parse_wps_via_zip(self, file_path: str) -> str:
        """
        尝试用 zipfile 直接解压 .wps 并提取 XML 中的文本。
        新版 WPS 本质是 ZIP 包，内含 word/document.xml 等。
        作为 python-docx 失败时的备选方案。
        """
        import zipfile
        from xml.etree import ElementTree as ET

        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                # 检查是否为有效的 ZIP (新版 WPS 必备特征)
                if 'word/document.xml' not in z.namelist():
                    return ""

                # 提取主文档 XML
                xml_content = z.read('word/document.xml')

                # 从 XML 中提取所有文本节点
                root = ET.fromstring(xml_content)
                ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                texts = []
                for t in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t'):
                    if t.text:
                        texts.append(t.text)
                return "\n".join(texts)
        except Exception:
            pass
        return ""

    def _parse_wps_via_wps_com(self, file_path: str) -> str:
        """
        通过 WPS Office COM 自动化提取文本。
        WPS Office 提供与 MS Word 兼容的 COM 接口 (ProgID: "WPS.Application"
        或 "KWPS.Application" 或 "ET.Application")。
        """
        abs_path = os.path.abspath(file_path)
        wps = None
        doc = None

        try:
            import pythoncom
            import win32com.client
        except ImportError:
            return ""

        # WPS Office 可能的 COM ProgID (按优先级)
        progids = [
            "WPS.Application",       # WPS Office 标准安装
            "KWPS.Application",      # Kingsoft WPS (旧版)
            "WPS.Document",          # 直接文档对象
            "Word.Application",      # MS Word (已在上层尝试)
        ]

        for progid in progids:
            if progid == "Word.Application":
                continue  # 已由上层 _parse_doc_via_win32 尝试

            try:
                pythoncom.CoInitialize()
                wps = win32com.client.Dispatch(progid)
                wps.Visible = False
                wps.DisplayAlerts = 0

                try:
                    doc = wps.Documents.Open(abs_path, ReadOnly=True, Visible=False)
                    text = doc.Content.Text

                    if text and text.strip():
                        return text
                except Exception:
                    pass
                finally:
                    if doc is not None:
                        try:
                            doc.Close(SaveChanges=False)
                        except Exception:
                            pass
                        doc = None
                    if wps is not None:
                        try:
                            wps.Quit()
                        except Exception:
                            pass
                        wps = None
            except Exception:
                pass
            finally:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

        return ""

    def _parse_ppt_file(self, file_path: str) -> str:
        """
        解析旧版 .ppt 文件 (PowerPoint 97-2003, OLE2 二进制)

        按优先级尝试:
          1. LibreOffice headless → txt (通用, 跨平台)
          2. win32com PowerPoint COM 自动化 (Windows + Office)
        """
        import subprocess
        import tempfile

        abs_path = os.path.abspath(file_path)

        # 方案1: LibreOffice headless 转 txt
        lo_paths = [
            "libreoffice", "soffice",
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        ]
        for lo_exe in lo_paths:
            try:
                subprocess.run([lo_exe, "--version"], capture_output=True, timeout=5)
                break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                lo_exe = None
                continue

        if lo_exe:
            try:
                with tempfile.TemporaryDirectory() as tmpdir:
                    cmd = [lo_exe, "--headless", "--convert-to", "txt:Text",
                           "--outdir", tmpdir, abs_path]
                    result = subprocess.run(cmd, capture_output=True, timeout=120)
                    if result.returncode == 0:
                        for f in os.listdir(tmpdir):
                            if f.endswith(".txt"):
                                txt_path = os.path.join(tmpdir, f)
                                with open(txt_path, "r", encoding="utf-8", errors="ignore") as fp:
                                    text = fp.read()
                                if text.strip():
                                    return text
            except Exception:
                pass

        # 方案2: PowerPoint COM 自动化
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            return ""

        ppt_app = None
        presentation = None
        try:
            pythoncom.CoInitialize()
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = False
            ppt_app.DisplayAlerts = 0
            presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True)

            slides_text = []
            for slide in presentation.Slides:
                texts = []
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        t = shape.TextFrame.TextRange.Text
                        if t and t.strip():
                            texts.append(t.strip())
                if texts:
                    slides_text.append("\n".join(texts))

            return "\n\n".join(slides_text)
        except Exception:
            return ""
        finally:
            if presentation is not None:
                try: presentation.Close()
                except Exception: pass
            if ppt_app is not None:
                try: ppt_app.Quit()
                except Exception: pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

    def _parse_ofd_file(self, file_path: str) -> str:
        """
        OFD 文件解析 — 带多级回退

        国产版式文档 OFD (GB/T 33190) 本质是 ZIP+XML 容器。
        优先级:
          1. 自定义提取器 (零依赖, 处理 GBK/UTF-8 混合编码 + 字体映射)
          2. ofdparser 库 (需安装: pip install ofdparser reportlab xmltodict)
        """
        # 方案1: 自定义提取器 (可靠, 无外部依赖)
        from ingestion.ofd_extractor import extract_ofd_text
        text = extract_ofd_text(file_path)
        if text and text.strip():
            return text

        # 方案2: ofdparser 库 (兜底, 处理特殊 OFD 变体)
        from ingestion.ofd_extractor import extract_ofd_text_via_ofdparser
        text = extract_ofd_text_via_ofdparser(file_path)
        if text and text.strip():
            return text

        return ""
